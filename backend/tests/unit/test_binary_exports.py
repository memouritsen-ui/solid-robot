"""Tests for binary export formats (PDF, DOCX, PPTX, XLSX).

Tasks #262-265: Test binary exports can be opened correctly.
"""

from io import BytesIO

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from research_tool.services.export import ExportFormat, ResearchExportData
from research_tool.services.export.docx import DOCXExporter
from research_tool.services.export.pdf import PDFExporter
from research_tool.services.export.pptx import PPTXExporter
from research_tool.services.export.xlsx import XLSXExporter


@pytest.fixture
def sample_research_data() -> ResearchExportData:
    """Create sample research data for testing."""
    return ResearchExportData(
        query="What are the effects of caffeine on sleep?",
        domain="medical",
        summary=(
            "Caffeine is a central nervous system stimulant that can "
            "significantly impact sleep quality and duration."
        ),
        facts=[
            {
                "statement": "Caffeine has a half-life of 5-6 hours",
                "confidence": 0.92,
                "source": "FDA",
                "verified": True,
            },
            {
                "statement": "Caffeine blocks adenosine receptors",
                "confidence": 0.95,
                "source": "PubMed",
                "verified": True,
            },
        ],
        sources=[
            {
                "title": "FDA Caffeine Guidelines",
                "url": "https://fda.gov/caffeine",
                "type": "government",
                "reliability_score": 0.95,
            },
            {
                "title": "Sleep Research Journal",
                "url": "https://sleepjournal.org/caffeine",
                "type": "academic",
                "reliability_score": 0.88,
            },
        ],
        confidence_score=0.90,
        limitations=[
            "Limited to adult populations",
            "Does not cover caffeine sensitivity variations",
        ],
        metadata={"research_id": "test-binary-export"},
    )


class TestPDFExport:
    """Tests for PDF export (#262)."""

    @pytest.fixture
    def exporter(self) -> PDFExporter:
        """Create PDFExporter instance."""
        return PDFExporter()

    def test_format_property(self, exporter: PDFExporter) -> None:
        """Test format property returns PDF."""
        assert exporter.format == ExportFormat.PDF

    def test_mime_type_property(self, exporter: PDFExporter) -> None:
        """Test mime_type property."""
        assert exporter.mime_type == "application/pdf"

    def test_file_extension_property(self, exporter: PDFExporter) -> None:
        """Test file_extension property."""
        assert exporter.file_extension == "pdf"

    @pytest.mark.asyncio
    async def test_export_returns_success(
        self, exporter: PDFExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns successful result."""
        result = await exporter.export(sample_research_data)

        assert result.success is True
        assert result.format == ExportFormat.PDF
        assert result.error is None

    @pytest.mark.asyncio
    async def test_export_returns_bytes(
        self, exporter: PDFExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns bytes content."""
        result = await exporter.export(sample_research_data)

        assert isinstance(result.content, bytes)
        assert len(result.content) > 0

    @pytest.mark.asyncio
    async def test_pdf_starts_with_magic_bytes(
        self, exporter: PDFExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test PDF content starts with PDF magic bytes."""
        result = await exporter.export(sample_research_data)

        # PDF files start with %PDF
        assert result.content[:4] == b"%PDF"

    @pytest.mark.asyncio
    async def test_pdf_is_valid_size(
        self, exporter: PDFExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test PDF has reasonable size (content was rendered)."""
        result = await exporter.export(sample_research_data)

        # PDF should be at least 1KB (has actual content)
        assert len(result.content) > 1000

    def test_generate_filename(self, exporter: PDFExporter) -> None:
        """Test filename generation."""
        filename = exporter.generate_filename("Test query")
        assert filename.endswith(".pdf")


class TestDOCXExport:
    """Tests for DOCX export (#263)."""

    @pytest.fixture
    def exporter(self) -> DOCXExporter:
        """Create DOCXExporter instance."""
        return DOCXExporter()

    def test_format_property(self, exporter: DOCXExporter) -> None:
        """Test format property returns DOCX."""
        assert exporter.format == ExportFormat.DOCX

    def test_mime_type_property(self, exporter: DOCXExporter) -> None:
        """Test mime_type property."""
        assert "wordprocessingml" in exporter.mime_type

    def test_file_extension_property(self, exporter: DOCXExporter) -> None:
        """Test file_extension property."""
        assert exporter.file_extension == "docx"

    @pytest.mark.asyncio
    async def test_export_returns_success(
        self, exporter: DOCXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns successful result."""
        result = await exporter.export(sample_research_data)

        assert result.success is True
        assert result.format == ExportFormat.DOCX
        assert result.error is None

    @pytest.mark.asyncio
    async def test_export_returns_bytes(
        self, exporter: DOCXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns bytes content."""
        result = await exporter.export(sample_research_data)

        assert isinstance(result.content, bytes)
        assert len(result.content) > 0

    @pytest.mark.asyncio
    async def test_docx_can_be_opened(
        self, exporter: DOCXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test DOCX can be opened with python-docx."""
        result = await exporter.export(sample_research_data)

        # Should not raise exception
        doc = Document(BytesIO(result.content))
        assert doc is not None

    @pytest.mark.asyncio
    async def test_docx_contains_title(
        self, exporter: DOCXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test DOCX contains the research title."""
        result = await exporter.export(sample_research_data)
        doc = Document(BytesIO(result.content))

        # Get all text from document
        full_text = "\n".join([p.text for p in doc.paragraphs])
        assert "Research Report" in full_text

    @pytest.mark.asyncio
    async def test_docx_contains_summary(
        self, exporter: DOCXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test DOCX contains the executive summary."""
        result = await exporter.export(sample_research_data)
        doc = Document(BytesIO(result.content))

        full_text = "\n".join([p.text for p in doc.paragraphs])
        assert "nervous system stimulant" in full_text

    @pytest.mark.asyncio
    async def test_docx_contains_findings(
        self, exporter: DOCXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test DOCX contains findings."""
        result = await exporter.export(sample_research_data)
        doc = Document(BytesIO(result.content))

        full_text = "\n".join([p.text for p in doc.paragraphs])
        assert "half-life" in full_text

    def test_generate_filename(self, exporter: DOCXExporter) -> None:
        """Test filename generation."""
        filename = exporter.generate_filename("Test query")
        assert filename.endswith(".docx")


class TestPPTXExport:
    """Tests for PPTX export (#264)."""

    @pytest.fixture
    def exporter(self) -> PPTXExporter:
        """Create PPTXExporter instance."""
        return PPTXExporter()

    def test_format_property(self, exporter: PPTXExporter) -> None:
        """Test format property returns PPTX."""
        assert exporter.format == ExportFormat.PPTX

    def test_mime_type_property(self, exporter: PPTXExporter) -> None:
        """Test mime_type property."""
        assert "presentationml" in exporter.mime_type

    def test_file_extension_property(self, exporter: PPTXExporter) -> None:
        """Test file_extension property."""
        assert exporter.file_extension == "pptx"

    @pytest.mark.asyncio
    async def test_export_returns_success(
        self, exporter: PPTXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns successful result."""
        result = await exporter.export(sample_research_data)

        assert result.success is True
        assert result.format == ExportFormat.PPTX
        assert result.error is None

    @pytest.mark.asyncio
    async def test_export_returns_bytes(
        self, exporter: PPTXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns bytes content."""
        result = await exporter.export(sample_research_data)

        assert isinstance(result.content, bytes)
        assert len(result.content) > 0

    @pytest.mark.asyncio
    async def test_pptx_can_be_opened(
        self, exporter: PPTXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test PPTX can be opened with python-pptx."""
        result = await exporter.export(sample_research_data)

        # Should not raise exception
        prs = Presentation(BytesIO(result.content))
        assert prs is not None

    @pytest.mark.asyncio
    async def test_pptx_has_slides(
        self, exporter: PPTXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test PPTX has multiple slides."""
        result = await exporter.export(sample_research_data)
        prs = Presentation(BytesIO(result.content))

        # Should have at least title, summary, findings, sources, limitations, closing
        assert len(prs.slides) >= 4

    @pytest.mark.asyncio
    async def test_pptx_title_slide(
        self, exporter: PPTXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test PPTX has proper title slide."""
        result = await exporter.export(sample_research_data)
        prs = Presentation(BytesIO(result.content))

        # First slide should be title slide
        title_slide = prs.slides[0]
        title_text = title_slide.shapes.title.text
        assert "Research Report" in title_text

    def test_generate_filename(self, exporter: PPTXExporter) -> None:
        """Test filename generation."""
        filename = exporter.generate_filename("Test query")
        assert filename.endswith(".pptx")


class TestXLSXExport:
    """Tests for XLSX export (#265)."""

    @pytest.fixture
    def exporter(self) -> XLSXExporter:
        """Create XLSXExporter instance."""
        return XLSXExporter()

    def test_format_property(self, exporter: XLSXExporter) -> None:
        """Test format property returns XLSX."""
        assert exporter.format == ExportFormat.XLSX

    def test_mime_type_property(self, exporter: XLSXExporter) -> None:
        """Test mime_type property."""
        assert "spreadsheetml" in exporter.mime_type

    def test_file_extension_property(self, exporter: XLSXExporter) -> None:
        """Test file_extension property."""
        assert exporter.file_extension == "xlsx"

    @pytest.mark.asyncio
    async def test_export_returns_success(
        self, exporter: XLSXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns successful result."""
        result = await exporter.export(sample_research_data)

        assert result.success is True
        assert result.format == ExportFormat.XLSX
        assert result.error is None

    @pytest.mark.asyncio
    async def test_export_returns_bytes(
        self, exporter: XLSXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test export returns bytes content."""
        result = await exporter.export(sample_research_data)

        assert isinstance(result.content, bytes)
        assert len(result.content) > 0

    @pytest.mark.asyncio
    async def test_xlsx_can_be_opened(
        self, exporter: XLSXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test XLSX can be opened with openpyxl."""
        result = await exporter.export(sample_research_data)

        # Should not raise exception
        wb = load_workbook(BytesIO(result.content))
        assert wb is not None

    @pytest.mark.asyncio
    async def test_xlsx_has_multiple_sheets(
        self, exporter: XLSXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test XLSX has multiple worksheets."""
        result = await exporter.export(sample_research_data)
        wb = load_workbook(BytesIO(result.content))

        # Should have Summary, Findings, Sources, Limitations sheets
        assert len(wb.sheetnames) >= 4
        assert "Summary" in wb.sheetnames
        assert "Findings" in wb.sheetnames
        assert "Sources" in wb.sheetnames

    @pytest.mark.asyncio
    async def test_xlsx_summary_sheet_contains_query(
        self, exporter: XLSXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test XLSX summary sheet contains query."""
        result = await exporter.export(sample_research_data)
        wb = load_workbook(BytesIO(result.content))

        ws = wb["Summary"]
        assert ws["B3"].value == sample_research_data.query

    @pytest.mark.asyncio
    async def test_xlsx_findings_sheet_has_data(
        self, exporter: XLSXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test XLSX findings sheet has fact data."""
        result = await exporter.export(sample_research_data)
        wb = load_workbook(BytesIO(result.content))

        ws = wb["Findings"]
        # Header row + 2 facts = 3 rows minimum
        assert ws.max_row >= 3
        # Check header
        assert ws["B1"].value == "Finding"

    @pytest.mark.asyncio
    async def test_xlsx_sources_sheet_has_data(
        self, exporter: XLSXExporter, sample_research_data: ResearchExportData
    ) -> None:
        """Test XLSX sources sheet has source data."""
        result = await exporter.export(sample_research_data)
        wb = load_workbook(BytesIO(result.content))

        ws = wb["Sources"]
        # Header row + 2 sources = 3 rows minimum
        assert ws.max_row >= 3

    def test_generate_filename(self, exporter: XLSXExporter) -> None:
        """Test filename generation."""
        filename = exporter.generate_filename("Test query")
        assert filename.endswith(".xlsx")
