"""PPTX export provider using python-pptx."""

from datetime import datetime
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from .exporter import Exporter, ExportFormat, ExportResult, ResearchExportData


class PPTXExporter(Exporter):
    """Export research results to PPTX format using python-pptx."""

    @property
    def format(self) -> ExportFormat:
        """Return the export format."""
        return ExportFormat.PPTX

    @property
    def mime_type(self) -> str:
        """Return the MIME type."""
        return (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        )

    @property
    def file_extension(self) -> str:
        """Return the file extension."""
        return "pptx"

    async def export(self, data: ResearchExportData) -> ExportResult:
        """Export research data to PPTX.

        Args:
            data: Research data to export

        Returns:
            ExportResult: PPTX content as bytes
        """
        try:
            prs = self._generate_presentation(data)
            buffer = BytesIO()
            prs.save(buffer)
            pptx_bytes = buffer.getvalue()

            return ExportResult(
                success=True,
                format=self.format,
                content=pptx_bytes,
                filename=self.generate_filename(data.query),
                mime_type=self.mime_type,
            )
        except Exception as e:
            return ExportResult(
                success=False,
                format=self.format,
                content=None,
                filename=self.generate_filename(data.query),
                mime_type=self.mime_type,
                error=str(e),
            )

    def _generate_presentation(self, data: ResearchExportData) -> Any:
        """Generate PPTX presentation from research data.

        Args:
            data: Research data

        Returns:
            Presentation: python-pptx Presentation object
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title Slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Research Report"
        subtitle.text = f"{data.query}\n\n{datetime.now().strftime('%Y-%m-%d')}"

        # Summary Slide
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "Executive Summary"

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = data.summary[:500]  # Limit for slide

        # Add metadata
        p = tf.add_paragraph()
        p.text = f"\nDomain: {data.domain}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"Confidence: {data.confidence_score:.1%}"
        p.level = 1

        # Key Findings Slides (max 3 findings per slide)
        if data.facts:
            for i in range(0, len(data.facts), 3):
                chunk = data.facts[i : i + 3]

                slide = prs.slides.add_slide(bullet_layout)
                slide.shapes.title.text = f"Key Findings ({i + 1}-{i + len(chunk)})"

                body = slide.placeholders[1]
                tf = body.text_frame

                for j, fact in enumerate(chunk):
                    statement = fact.get("statement", fact.get("content", ""))
                    confidence = fact.get("confidence", 0.0)

                    if j == 0:
                        tf.text = statement[:200]
                    else:
                        p = tf.add_paragraph()
                        p.text = statement[:200]
                        p.level = 0

                    p = tf.add_paragraph()
                    p.text = f"Confidence: {confidence:.1%}"
                    p.level = 1

        # Sources Slide
        if data.sources:
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = "Sources"

            body = slide.placeholders[1]
            tf = body.text_frame

            for i, source in enumerate(data.sources[:8]):  # Max 8 sources on slide
                title = source.get("title", "Untitled")

                if i == 0:
                    tf.text = title
                else:
                    p = tf.add_paragraph()
                    p.text = title
                    p.level = 0

            if len(data.sources) > 8:
                p = tf.add_paragraph()
                p.text = f"... and {len(data.sources) - 8} more sources"
                p.level = 1

        # Limitations Slide
        if data.limitations:
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = "Limitations"

            body = slide.placeholders[1]
            tf = body.text_frame

            for i, limitation in enumerate(data.limitations[:6]):
                if i == 0:
                    tf.text = limitation
                else:
                    p = tf.add_paragraph()
                    p.text = limitation
                    p.level = 0

        # Closing Slide
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = "Questions?"
        slide.placeholders[1].text = (
            "This report was generated automatically.\n"
            "Please verify critical information from primary sources."
        )

        return prs
